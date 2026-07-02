from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

def create_presentation_deck():
    output_dir = r"C:\Users\rkvig\Downloads\Joshi\task4_storytelling"
    output_path = os.path.join(output_dir, "final_presentation.pptx")
    
    os.makedirs(output_dir, exist_ok=True)
    
    prs = Presentation()
    
    # Standard 16:9 widescreen layout
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Custom colors
    TEAL = RGBColor(15, 76, 92)      # Primary theme
    NAVY = RGBColor(29, 53, 87)      # Text/Headers
    ORANGE = RGBColor(227, 100, 20)  # Accent
    SLATE = RGBColor(100, 116, 139)  # Subtitles/Muted
    WHITE = RGBColor(255, 255, 255)
    
    # Helper to set background color of a slide
    def set_slide_bg(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color
        
    # Helper to add a formatted title to a slide
    def add_slide_header(slide, title_text, category_text="DATA ANALYTICS INTERNSHIP"):
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.83), Inches(1.2))
        tf = title_box.text_frame
        tf.word_wrap = True
        
        # Category label
        p_cat = tf.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.name = 'Arial'
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = ORANGE
        p_cat.space_after = Pt(2)
        
        # Main Title
        p_title = tf.add_paragraph()
        p_title.text = title_text
        p_title.font.name = 'Arial'
        p_title.font.size = Pt(24)
        p_title.font.bold = True
        p_title.font.color.rgb = NAVY
        
    # -------------------------------------------------------------
    # SLIDE 1: Title Slide (Dark Background)
    # -------------------------------------------------------------
    slide_layout = prs.slide_layouts[6] # Blank slide
    slide1 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide1, TEAL)
    
    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.33), Inches(3.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p_sub = tf.paragraphs[0]
    p_sub.text = "APEXPLANET SOFTWARE PVT. LTD. - PORTFOLIO PRESENTATION"
    p_sub.font.name = 'Arial'
    p_sub.font.size = Pt(12)
    p_sub.font.bold = True
    p_sub.font.color.rgb = ORANGE
    p_sub.space_after = Pt(10)
    
    p_main = tf.add_paragraph()
    p_main.text = "Online Retail Sales & Customer Segmentation"
    p_main.font.name = 'Arial'
    p_main.font.size = Pt(40)
    p_main.font.bold = True
    p_main.font.color.rgb = WHITE
    p_main.space_after = Pt(10)
    
    p_author = tf.add_paragraph()
    p_author.text = "Prepared by: Data Analytics Intern | Final Capstone Project"
    p_author.font.name = 'Arial'
    p_author.font.size = Pt(14)
    p_author.font.italic = True
    p_author.font.color.rgb = WHITE
    
    # -------------------------------------------------------------
    # SLIDE 2: Executive Summary
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(slide_layout)
    add_slide_header(slide2, "Executive Overview & Core Metrics")
    
    # Content block
    content_box = slide2.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.83), Inches(5.0))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Over the 12-month period analyzed, the retail store generated robust sales driven primarily by a loyal domestic buyer base and high-value B2B international accounts."
    p.font.size = Pt(15)
    p.font.name = 'Arial'
    p.space_after = Pt(15)
    
    # Bullet points of metrics
    bullets = [
        "Total Net Revenue: GBP 9.74 Million (GBP 9,742,030.56)",
        "Total Shipped Orders: 19,960 completed transactions",
        "Unique Active Buyers: 4,321 unique customers (registered and guest checkouts)",
        "Average Order Value (AOV): GBP 488.08 per transaction basket",
        "Repeat Buyer Retention: 65.58% repeat customers, indicating strong brand loyalty"
    ]
    
    for b in bullets:
        p_b = tf.add_paragraph()
        p_b.text = "• " + b
        p_b.font.size = Pt(14)
        p_b.font.name = 'Arial'
        p_b.space_after = Pt(10)
        p_b.level = 0
        
    # -------------------------------------------------------------
    # SLIDE 3: Task 1 - Data Wrangling
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(slide_layout)
    add_slide_header(slide3, "Task 1: Data Wrangling & Quality Cleaning")
    
    content_box = slide3.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.83), Inches(5.0))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    bullets_t1 = [
        "Encoding Normalization: Resolved non-UTF-8 characters in descriptions using ISO-8859-1 mapping.",
        "Missing Data Handling: Approximately 25% of transactions lacked a CustomerID. Guest transactions were isolated to 'Guest_<InvoiceNo>' to preserve net sales statistics while avoiding segment distortion.",
        "Duplicate Elimination: Dropped 5,268 double-entered rows (0.97% of database).",
        "Transaction Anomaly Filtering: Excluded non-cancellation negative/zero quantity rows representing damage write-offs, and removed zero-price items to maintain financial integrity.",
        "Output: Compiled a clean, indexed dataset of 534,129 rows (98.56% of raw database)."
    ]
    
    tf.paragraphs[0].text = "Cleaning raw transactional logs into analysis-ready databases:"
    tf.paragraphs[0].font.size = Pt(15)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].space_after = Pt(15)
    
    for b in bullets_t1:
        p_b = tf.add_paragraph()
        p_b.text = "• " + b
        p_b.font.size = Pt(14)
        p_b.space_after = Pt(12)
        
    # -------------------------------------------------------------
    # SLIDE 4: Task 2 - Seasonal & Operational Trends
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(slide_layout)
    add_slide_header(slide4, "Task 2: Seasonality & Operational Insights")
    
    content_box = slide4.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.83), Inches(5.0))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    bullets_t2 = [
        "Strong Holiday Seasonality: Monthly sales show a massive autumn spike. November 2011 peaked at GBP 1.46 Million (15% of annual sales), driven by holiday stocking.",
        "Domestic vs. Export: The United Kingdom represents 92.4% of total net revenue. The top international export market is the Netherlands, generating GBP 285k in revenue.",
        "Midday Order Density: Transactions are highly concentrated during working hours (10:00 AM - 3:00 PM), peaking at 12:00 PM. Operational picking queues should be staffed accordingly.",
        "Constant Basket Size: Despite massive fluctuations in order volume throughout the day, item counts per basket remain constant, suggesting buying patterns are stable."
    ]
    
    tf.paragraphs[0].text = "Uncovering critical revenue and purchase patterns:"
    tf.paragraphs[0].font.size = Pt(15)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].space_after = Pt(15)
    
    for b in bullets_t2:
        p_b = tf.add_paragraph()
        p_b.text = "• " + b
        p_b.font.size = Pt(14)
        p_b.space_after = Pt(12)
        
    # -------------------------------------------------------------
    # SLIDE 5: Task 3 - Customer Segmentation (RFM)
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(slide_layout)
    add_slide_header(slide5, "Task 3: RFM Customer Cohort Segmentation")
    
    content_box = slide5.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.83), Inches(5.0))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    bullets_t3 = [
        "Loyal Customers (26.96%): Buy frequently and recently. Average spend is GBP 1,911.39.",
        "Champions (18.61%): Top tier. Average spend is GBP 6,777.62 with a recency of 15 days.",
        "About to Sleep (19.83%): Below average recency and frequency. Target of win-back marketing.",
        "Lost (22.70%): Long-term inactive, single purchase. Low priority re-engagement cohort.",
        "At Risk (4.05%): High spending, but inactive for 250+ days. Needs immediate win-back campaigns."
    ]
    
    tf.paragraphs[0].text = "Partitioning 4,321 registered buyers into behavioral cohorts:"
    tf.paragraphs[0].font.size = Pt(15)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].space_after = Pt(15)
    
    for b in bullets_t3:
        p_b = tf.add_paragraph()
        p_b.text = "• " + b
        p_b.font.size = Pt(14)
        p_b.space_after = Pt(10)
        
    # -------------------------------------------------------------
    # SLIDE 6: Task 4 - Hypothesis Testing
    # -------------------------------------------------------------
    slide6 = prs.slides.add_slide(slide_layout)
    add_slide_header(slide6, "Task 4: Statistical Hypothesis Validation")
    
    content_box = slide6.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.83), Inches(5.0))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    bullets_t4 = [
        "Test 1: UK vs. International Average Order Value (AOV)",
        "  - Null Hypothesis (H0): UK AOV = International AOV",
        "  - Result: Welch's T-Test yielded T-statistic = -8.29, P-value = 1.8e-16 (Highly Significant).",
        "  - Decision: Reject H0. International orders have a significantly larger average size (GBP 845.11 vs GBP 499.57).",
        "Test 2: Peak vs. Off-Peak Order Size",
        "  - Null Hypothesis (H0): Mean items per order are identical during peak (10 AM - 3 PM) and off-peak.",
        "  - Result: Welch's T-Test yielded T-statistic = -1.14, P-value = 0.255 (Not Significant).",
        "  - Decision: Fail to reject H0. Order sizes are statistically identical throughout operating hours."
    ]
    
    tf.paragraphs[0].text = "Applying statistical rigor to validate business patterns:"
    tf.paragraphs[0].font.size = Pt(15)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].space_after = Pt(15)
    
    for b in bullets_t4:
        p_b = tf.add_paragraph()
        p_b.text = b
        p_b.font.size = Pt(13)
        p_b.space_after = Pt(6)
        
    # -------------------------------------------------------------
    # SLIDE 7: Strategic Recommendations
    # -------------------------------------------------------------
    slide7 = prs.slides.add_slide(slide_layout)
    add_slide_header(slide7, "Operational & Strategic Action Plan")
    
    content_box = slide7.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.83), Inches(5.0))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    bullets_t5 = [
        "B2B Shipping Optimization: Introduce bulk shipping logistics tiers for top international zones (Netherlands, Australia, EIRE) to reward and grow high-AOV wholesale customers.",
        "Targeted Reactivation Campaigns: Deploy automated emails with free shipping/promo offers to the 175 high-value 'At Risk' customers before they churn completely.",
        "Hour-based Staff Scheduling: Adjust warehouse picking staff schedules to match peak incoming transaction density (midday peak window: 10:00 AM - 3:00 PM).",
        "New Customer Welcome Funnels: Offer a 'second-purchase incentive' within 14 days of checkout to the 205 Recent/New Customers to nurture them into Loyal Customers."
    ]
    
    tf.paragraphs[0].text = "Key business actions based on data insights:"
    tf.paragraphs[0].font.size = Pt(15)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].space_after = Pt(15)
    
    for b in bullets_t5:
        p_b = tf.add_paragraph()
        p_b.text = "• " + b
        p_b.font.size = Pt(14)
        p_b.space_after = Pt(12)
        
    prs.save(output_path)
    print(f"Presentation deck successfully saved to {output_path}")

if __name__ == "__main__":
    create_presentation_deck()